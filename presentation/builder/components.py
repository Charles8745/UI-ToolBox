"""元件庫頁：KPI 卡、note bar、chip、編號卡、章節 nav 膠囊、styled 表格。全原生形狀。"""
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from . import PX
from .shapes import add_accent_bar, add_glass_card, add_text


def _blank_on(prs, layouts, name):
    slide = prs.slides.add_slide(layouts[name])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def add_component_slides(prs, layouts, tokens) -> None:
    t = tokens
    # --- 元件庫 1：KPI 卡 / note bar / chips / nav 膠囊 ---
    s1 = _blank_on(prs, layouts, "LG 內容")
    add_text(s1.shapes, 96, 40, 900, 50, "元件庫 1 / KPI · note bar · chip · nav", 20, t["textDim"], tokens=t)

    # KPI 卡（玻璃卡 + mono 大數字 + 標籤）
    add_glass_card(s1.shapes, 96, 130, 470, 270, t)
    add_text(s1.shapes, 128, 170, 400, 90, "1,728 萬", 40, t["accent1"], bold=True, mono=True, tokens=t)
    add_text(s1.shapes, 128, 290, 400, 50, "KPI 標籤文字", 16, t["textDim"], tokens=t)

    # note bar（玻璃長條 + accent 左標）
    add_glass_card(s1.shapes, 96, 460, 1728, 110, t)
    add_accent_bar(s1.shapes, 120, 484, 8, 62, t["accent1"])
    add_text(s1.shapes, 156, 486, 120, 50, "結論", 18, t["accent1"], bold=True, tokens=t)
    add_text(s1.shapes, 290, 486, 1480, 60, "note bar 說明文字：一行重點結論或補充。", 18, t["text"], tokens=t)

    # chips（三個小膠囊）
    for i, (label, color) in enumerate([("LIVE", t["accent1"]), ("MOCK", t["accent3"]), ("ALERT", t["accent4"])]):
        x = 96 + i * 220
        card = add_glass_card(s1.shapes, x, 640, 180, 64, t)
        card.adjustments[0] = 0.5
        add_text(s1.shapes, x, 652, 180, 40, label, 14, color, bold=True, align="c", tokens=t)

    # 章節 nav 膠囊列（六格）
    for i in range(6):
        x = 96 + i * 300
        card = add_glass_card(s1.shapes, x, 780, 280, 70, t)
        card.adjustments[0] = 0.5
        add_text(s1.shapes, x, 796, 280, 40, f"Section {i + 1}", 14, t["text"], align="c", tokens=t)

    # --- 元件庫 2：編號卡 / styled 表格 ---
    s2 = _blank_on(prs, layouts, "LG 內容")
    add_text(s2.shapes, 96, 40, 900, 50, "元件庫 2 / 編號卡 · 表格", 20, t["textDim"], tokens=t)

    add_glass_card(s2.shapes, 96, 130, 420, 420, t)
    add_text(s2.shapes, 128, 160, 200, 90, "01", 40, t["accent4"], bold=True, mono=True, tokens=t)
    add_text(s2.shapes, 128, 280, 360, 60, "編號卡標題", 22, t["text"], bold=True, tokens=t)
    add_text(s2.shapes, 128, 350, 360, 140, "編號卡描述文字，兩到三行的說明內容。", 16, t["textDim"], tokens=t)

    # styled 表格：表頭 surface 填色 + accent 字、條紋列
    rows, cols = 4, 3
    gf = s2.shapes.add_table(rows, cols, Emu(600 * PX), Emu(130 * PX), Emu(1220 * PX), Emu(420 * PX))
    table = gf.table
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = f"欄位{c + 1}"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(t["surface"])
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.color.rgb = RGBColor.from_string(t["accent1"])
        para.runs[0].font.bold = True
        para.runs[0].font.size = Pt(16)
    for r in range(1, rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = "內容"
            cell.fill.solid()
            shade = "10141B" if r % 2 else "0B111A"
            cell.fill.fore_color.rgb = RGBColor.from_string(shade)
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.color.rgb = RGBColor.from_string(t["text"])
            para.runs[0].font.size = Pt(14)
