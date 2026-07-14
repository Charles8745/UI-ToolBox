"""原生玻璃形狀工廠：圓角矩形 + 半透明填色 + 髮絲線 + 柔影。座標一律 px。"""
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from . import PX

A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}


def _append_alpha(srgb_el, pct: int) -> None:
    srgb_el.append(srgb_el.makeelement(qn("a:alpha"), {"val": str(pct * 1000)}))


def _fill_srgb(shape):
    # 走 spPr 的 XML 路徑，不碰 ColorFormat 私有屬性
    return shape._element.spPr.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))


def add_glass_card(shapes, x, y, w, h, tokens):
    card = tokens["card"]
    sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Emu(x * PX), Emu(y * PX), Emu(w * PX), Emu(h * PX))
    sp.adjustments[0] = card["radius"]
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(card["fill"])
    _append_alpha(_fill_srgb(sp), card["fillAlpha"])

    sp.line.color.rgb = RGBColor.from_string("FFFFFF")
    sp.line.width = Pt(1)
    ln = sp._element.spPr.find(qn("a:ln"))
    _append_alpha(ln.find(".//" + qn("a:srgbClr")), card["lineAlpha"])

    # 注意：不可先 sp.shadow.inherit = False（會建空 effectLst，再 append 會出現兩個 effectLst，schema 違規）。
    # 直接 append 自己的 effectLst 即同時覆蓋繼承。
    effect = etree.fromstring(
        f"<a:effectLst {A}>"
        '<a:outerShdw blurRad="228600" dist="76200" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="35000"/></a:srgbClr>'
        "</a:outerShdw></a:effectLst>"
    )
    sp._element.spPr.append(effect)
    sp.text_frame.paragraphs[0].text = ""
    return sp


def add_text(shapes, x, y, w, h, text, sz, color, bold=False, mono=False, align="l", tokens=None):
    box = shapes.add_textbox(Emu(x * PX), Emu(y * PX), Emu(w * PX), Emu(h * PX))
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = _ALIGN[align]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = tokens["fontMono"] if mono else tokens["fontLatin"]
    if not mono:
        # East Asian 字型
        rpr = run._r.get_or_add_rPr()
        ea = rpr.makeelement(qn("a:ea"), {"typeface": tokens["fontEastAsian"]})
        rpr.append(ea)
    return box


def add_accent_bar(shapes, x, y, w, h, hex_color):
    sp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Emu(x * PX), Emu(y * PX), Emu(w * PX), Emu(h * PX))
    sp.adjustments[0] = 0.5
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp
