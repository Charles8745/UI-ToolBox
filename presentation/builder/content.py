"""範例頁：每版型一頁假內容，最底層埋動態 mp4（autoplay + loop）。"""
from pptx.oxml.ns import qn
from pptx.util import Emu

from . import PX, ROOT
from .shapes import add_accent_bar, add_glass_card, add_text

MP4 = ROOT / "assets" / "media" / "wave-loop.mp4"
POSTER = ROOT / "assets" / "media" / "wave-frame.png"


def add_bg_movie(slide, prs) -> None:
    movie = slide.shapes.add_movie(
        str(MP4), 0, 0, prs.slide_width, prs.slide_height,
        poster_frame_image=str(POSTER), mime_type="video/mp4")
    # 移到最底層（spTree 前兩個子節點是 nvGrpSpPr/grpSpPr）
    el = movie._element
    tree = el.getparent()
    tree.remove(el)
    tree.insert(2, el)
    # autoplay：click 觸發的 delay=indefinite 改 0
    timing = slide._element.find(qn("p:timing"))
    for cond in timing.iter(qn("p:cond")):
        if cond.get("delay") == "indefinite":
            cond.set("delay", "0")
    # loop：cMediaNode 的 cTn 設 repeatCount
    for cmed in timing.iter(qn("p:cMediaNode")):
        ctn = cmed.find(qn("p:cTn"))
        if ctn is not None:
            ctn.set("repeatCount", "indefinite")


def _fill_placeholders(slide, texts: dict) -> None:
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in texts:
            ph.text = texts[idx]


def add_example_slides(prs, layouts, tokens) -> None:
    t = tokens
    # 每版型一頁；title 的 idx 為 0
    plan = [
        ("LG 封面", {0: "永續智能航港生態系", 1: "Liquid Glass 深色簡報模板 範例", 2: "團隊名稱 · 2026"}),
        ("LG Outline", {0: "OUTLINE"}),
        ("LG 流程", {0: "研究流程 Pipeline", 1: "OVERVIEW / 研究脈絡"}),
        ("LG 內容", {0: "標準內容頁標題", 1: "01 / MOTIVATION",
                     2: "左欄內容文字範例：\n重點一說明\n重點二說明\n重點三說明"}),
        ("LG 統計", {0: "資料集統計", 1: "02 / DATASET", 2: "主視覺圖說範例"}),
        ("LG 表格", {0: "結果表格", 1: "05 / RESULT"}),
        ("LG 比較", {0: "方法比較", 1: "04 / CLUSTERING"}),
        ("LG 氛圍", {0: "全幅氛圍頁標題範例"}),
        ("LG 結語", {0: "CONCLUSION", 1: "Thank You", 2: "團隊名單範例"}),
    ]
    for name, texts in plan:
        slide = prs.slides.add_slide(layouts[name])
        _fill_placeholders(slide, texts)
        add_bg_movie(slide, prs)

        if name == "LG 統計":
            # 範例 KPI 卡兩張
            add_glass_card(slide.shapes, 96, 300, 470, 270, t)
            add_text(slide.shapes, 128, 340, 400, 90, "100,000", 40, t["accent1"], bold=True, mono=True, tokens=t)
            add_text(slide.shapes, 128, 460, 400, 50, "KPI 標籤範例", 16, t["textDim"], tokens=t)
            add_glass_card(slide.shapes, 590, 300, 470, 270, t)
        if name == "LG 流程":
            for i in range(5):
                add_glass_card(slide.shapes, 96 + i * 356, 300, 320, 480, t)
            add_glass_card(slide.shapes, 96, 860, 1728, 110, t)
            add_accent_bar(slide.shapes, 120, 884, 8, 62, t["accent1"])
        if name == "LG 比較":
            add_glass_card(slide.shapes, 96, 800, 1728, 140, t)
            add_accent_bar(slide.shapes, 120, 824, 8, 62, t["accent1"])


def add_calibration_slide(prs, layouts, tokens) -> None:
    """--calibration：烘圖卡 vs 原生卡並排。"""
    slide = prs.slides.add_slide(layouts["LG 內容"])
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    baked = ROOT / "assets" / "calibration" / "card.png"
    slide.shapes.add_picture(str(baked), 0, 0, prs.slide_width, prs.slide_height)
    add_glass_card(slide.shapes, 1150, 360, 560, 360, tokens)
    add_text(slide.shapes, 1182, 420, 480, 90, "1,728 萬次", 40, tokens["accent1"], bold=True, mono=True, tokens=tokens)
    add_text(slide.shapes, 1182, 540, 480, 50, "原生形狀卡（比對右側烘圖）", 16, tokens["textDim"], tokens=tokens)
