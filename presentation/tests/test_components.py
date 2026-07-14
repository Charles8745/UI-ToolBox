import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))


@pytest.fixture(scope="module")
def prs():
    from builder import load_tokens
    from builder.theme import write_theme
    from builder.layouts import build_layouts
    from builder.components import add_component_slides

    p = Presentation(ROOT / "skeleton.pptx")
    tokens = load_tokens()
    write_theme(p, tokens)
    layouts = build_layouts(p, tokens)
    add_component_slides(p, layouts, tokens)
    return p


def test_component_slides_exist(prs):
    assert len(list(prs.slides)) >= 2


def test_glass_card_has_alpha_fill(prs):
    slide = list(prs.slides)[0]
    alphas = slide.shapes._spTree.findall(
        ".//" + qn("a:solidFill") + "/" + qn("a:srgbClr") + "/" + qn("a:alpha"))
    assert alphas, "元件庫頁找不到帶 alpha 的玻璃卡填色"
    assert any(a.get("val") == "55000" for a in alphas)


def test_kpi_number_uses_mono_font(prs):
    slide = list(prs.slides)[0]
    latins = [r.get("typeface") for r in slide.shapes._spTree.iter(qn("a:latin"))]
    assert "Geist Mono" in latins
