import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))


@pytest.fixture(scope="module")
def dist(tmp_path_factory):
    sys.path.insert(0, str(ROOT))
    import build_pptx

    out = tmp_path_factory.mktemp("dist") / "t.pptx"
    build_pptx.build(out, calibration=False)
    return out


def test_build_produces_file(dist):
    assert dist.exists() and dist.stat().st_size > 0


def test_example_slides_use_all_nine_layouts(dist):
    prs = Presentation(dist)
    used = {s.slide_layout.name for s in prs.slides}
    for n in ["LG 封面", "LG Outline", "LG 流程", "LG 內容", "LG 統計",
              "LG 表格", "LG 比較", "LG 氛圍", "LG 結語"]:
        assert n in used, f"範例頁缺 {n}"


def test_movie_embedded_and_autoplay(dist):
    prs = Presentation(dist)
    cover = next(s for s in prs.slides if s.slide_layout.name == "LG 封面")
    videos = cover._element.findall(".//" + qn("a:videoFile"))
    assert videos, "封面範例頁沒有影片"
    timing = cover._element.find(qn("p:timing"))
    conds = [c.get("delay") for c in timing.iter(qn("p:cond"))]
    assert "indefinite" not in conds, "影片仍是點擊播放（delay=indefinite 未改 0）"
    ctns = [c.get("repeatCount") for c in timing.iter(qn("p:cTn"))]
    assert "indefinite" in ctns, "影片未設 loop"


def test_file_size_sane(dist):
    assert dist.stat().st_size < 60 * 1024 * 1024, "產物超過 60MB，媒體疑似未去重"
