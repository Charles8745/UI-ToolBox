import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

EXPECTED = ["LG 封面", "LG Outline", "LG 流程", "LG 內容", "LG 統計",
            "LG 表格", "LG 比較", "LG 氛圍", "LG 結語"]


@pytest.fixture(scope="module")
def prs():
    from builder import load_tokens
    from builder.theme import write_theme
    from builder.layouts import build_layouts

    p = Presentation(ROOT / "skeleton.pptx")
    write_theme(p, load_tokens())
    build_layouts(p, load_tokens())
    return p


def test_nine_layouts_created(prs):
    names = [l.name for l in prs.slide_masters[0].slide_layouts]
    for n in EXPECTED:
        assert n in names, f"缺 layout {n}"


def test_layout_backgrounds_are_pictures(prs):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name not in EXPECTED:
            continue
        bg = layout.element.find(qn("p:cSld")).find(qn("p:bg"))
        assert bg is not None, f"{layout.name} 無背景"
        blip = bg.find(qn("p:bgPr")).find(qn("a:blipFill")).find(qn("a:blip"))
        r_id = blip.get(qn("r:embed"))
        part = layout.part.related_part(r_id)
        assert part.content_type.startswith("image/"), f"{layout.name} 背景非圖片"


def test_cover_placeholder_geometry(prs):
    cover = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "LG 封面")
    title = next(ph for ph in cover.placeholders if ph.placeholder_format.type is not None
                 and ph.placeholder_format.idx == 0)
    assert title.left == 180 * 6350
    assert title.top == 360 * 6350


def test_roundtrip_save_load(prs, tmp_path):
    out = tmp_path / "t.pptx"
    prs.save(out)
    p2 = Presentation(out)
    assert len(list(p2.slide_masters[0].slide_layouts)) >= 9
